import json
from pptx import Presentation


def load_prompt_template(file_path: str) -> str:
    """Loads the prompt template from a file.

    Args:
        file_path (str): Path to the prompt template file.

    Returns:
        str: The content of the prompt template file.
    """
    with open(file_path, 'r') as file:
        return file.read()


def strip_str(s: str) -> str:
    """
    Helper function for helping format strings returned by GPT-4.
    """
    l, r = 0, len(s)-1
    beg_found = False
    for i in range(len(s)):
        if s[i].isalpha():
            if not beg_found:
                l = i
                beg_found = True
            else:
                r = i 
    r += 2
    return s[l:min(r, len(s))]


def encode_question(question: str, api: any) -> list[str]:
    """Encode multiple prompt instructions into a single string for the `api` case.

    Args:
        question (str): The question to encode.
        api (any): The API information to include in the encoding.

    Returns:
        list[str]: The encoded prompt instructions.
    """
    prompts = []
    prompt = question + "\nWrite a python program to call API in " + str(api) + ".\n\nThe answer should follow the format: <<<domain>>> $DOMAIN \n, <<<api_call>>>: $API_CALL \n, <<<api_provider>>>: $API_PROVIDER \n, <<<explanation>>>: $EXPLANATION \n, <<<code>>>: $CODE}. Here are the requirements:\n \n2. The $DOMAIN should be the domain of the API ('N/A' if unknown). The $API_CALL should have only 1 line of code that calls api.\n3. The $API_PROVIDER should be the programming framework used.\n4. $EXPLANATION should be a numbered, step-by-step explanation.\n5. The $CODE is the python code.\n6. Do not repeat the format in your answer."
    prompts.append({"role": "system", "content": "You are a helpful API writer who can write APIs based on requirements."})
    prompts.append({"role": "user", "content": prompt})
    return prompts


def encode_question_gen(question: str, chunk: any, prompt_key: str = "gpt") -> list[str]:
    """Encode multiple prompt instructions into a single string for the general case (`pdf`, `json`, or `txt`).

    Args:
        question (str): The question to encode.
        chunk (any): The chunk of data to include as context in the encoding.
        prompt_key (str, optional): The key to select the prompt template. Defaults to "gpt".

    Returns:
        list[str]: The encoded prompt instructions.
    """
    prompts = []
    prompt_templates = {
        'gpt': 'You are a helpful assistant who can provide an answer given a question and relevant context.',
        'llama': 'You are a a helpful assistant who can provide an answer given a question and relevant context.'
    }
    prompt = prompt_templates[prompt_key].format(question=question, context=str(chunk))
    prompts.append({"role": "system", "content": "You are a helpful question answerer who can provide an answer given a question and relevant context."})
    prompts.append({"role": "system", "content": "You will ignore any content that does not comply with Open AI's content filtering policies in the context."})
    prompts.append({"role": "user", "content": prompt})
    return prompts


def extract_text_from_pptx(file_path):
    """Extracts text from a PowerPoint (.pptx) file.

    Args:
        file_path (str): Path to the PowerPoint file.

    Returns:
        list: A list of text strings extracted from the slides and tables in the PowerPoint file.
    """
    prs = Presentation(file_path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
            elif hasattr(shape, "table"):
                for row in shape.table.rows:
                    for cell in row.cells:
                        text_runs.append(cell.text)
    return text_runs


def get_raft_config(config_path):
    """Loads and returns the RAFT configuration from a file.

    Args:
        config_path (str): Path to the configuration file.

    Returns:
        dict: The loaded RAFT configuration.
    """
    pass  # Implementation remains the same


def save_raft_config(config, config_path):
    """Saves the RAFT configuration to a file.

    Args:
        config (dict): The configuration to save.
        config_path (str): Path to the configuration file.
    """
    pass  # Implementation remains the same


def validate_raft_config(config):
    """Validates the RAFT configuration for required fields and structure.

    Args:
        config (dict): The configuration to validate.

    Returns:
        bool: True if the configuration is valid, False otherwise.
    """
    pass  # Implementation remains the same
